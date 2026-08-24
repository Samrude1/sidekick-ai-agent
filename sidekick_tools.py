from playwright.async_api import async_playwright
from langchain_community.agent_toolkits import PlayWrightBrowserToolkit
from dotenv import load_dotenv
import os
import requests
import ast
import io
import contextlib
import math
import statistics
import random
import datetime as dt_module
from datetime import datetime
import json
import re
import collections
import itertools
import ipaddress
import socket
import urllib.parse
from langchain_core.tools import Tool, StructuredTool
from pydantic import BaseModel, Field
from langchain_community.tools.wikipedia.tool import WikipediaQueryRun
from langchain_community.utilities import SerpAPIWrapper
from langchain_community.utilities.wikipedia import WikipediaAPIWrapper
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load environment from root folder
load_dotenv(".env", override=True)

pushover_token = os.getenv("PUSHOVER_TOKEN")
pushover_user = os.getenv("PUSHOVER_USER")
pushover_url = "https://api.pushover.net/1/messages.json"

# Using SerpAPI (as per Lab 2 fix)
serp = SerpAPIWrapper()


def is_safe_url(url: str) -> tuple[bool, str]:
    """Validate URLs against SSRF, loopback, private networks, and cloud metadata endpoints."""
    try:
        url_str = url.strip()
        parsed = urllib.parse.urlparse(url_str)
        if parsed.scheme not in ("http", "https"):
            return False, f"Blocked: Protocol '{parsed.scheme}' is not permitted. Only 'http' and 'https' are allowed."
        
        hostname = parsed.hostname
        if not hostname:
            return False, "Blocked: Missing hostname in URL."
        
        hostname_lower = hostname.lower()
        if hostname_lower in ("localhost", "127.0.0.1", "0.0.0.0", "::1") or hostname_lower.endswith(".localhost") or hostname_lower.endswith(".local"):
            return False, f"Blocked: Navigation to local host '{hostname}' is forbidden for security."
        
        # Check IP literal or resolve hostname to IP
        try:
            ip_obj = ipaddress.ip_address(hostname_lower)
            ips = [ip_obj]
        except ValueError:
            try:
                addr_info = socket.getaddrinfo(hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM)
                ips = [ipaddress.ip_address(item[4][0]) for item in addr_info]
            except Exception as e:
                return False, f"Blocked: Unable to resolve hostname '{hostname}' ({e})."
        
        for ip in ips:
            if ip.is_loopback or ip.is_private or ip.is_link_local or ip.is_reserved or ip.is_unspecified:
                return False, f"Blocked: Target resolves to private/restricted IP address {ip}."
        
        return True, ""
    except Exception as e:
        return False, f"Blocked: Invalid URL format ({e})."


async def playwright_tools():
    playwright = await async_playwright().start()
    browser = await playwright.chromium.launch(headless=True)
    toolkit = PlayWrightBrowserToolkit.from_browser(async_browser=browser)
    tools = toolkit.get_tools()

    # Wrap navigation tools with SSRF validation
    for tool in tools:
        if tool.name == "navigate_browser":
            original_arun = tool._arun
            original_run = tool._run

            async def safe_navigate_arun(url: str, **kwargs):
                safe, reason = is_safe_url(url)
                if not safe:
                    return reason
                return await original_arun(url=url, **kwargs)

            def safe_navigate_run(url: str, **kwargs):
                safe, reason = is_safe_url(url)
                if not safe:
                    return reason
                return original_run(url=url, **kwargs)

            tool._arun = safe_navigate_arun
            tool._run = safe_navigate_run

    return tools, browser, playwright


push_count = 0

def push(text: str) -> str:
    """Send a push notification to the user via Pushover with rate limiting and timeout."""
    global push_count
    if push_count >= 2:
        return "Push limit reached for this session."
    push_count += 1
    
    if not pushover_token or not pushover_user:
        return "Push notification skipped: Pushover credentials not configured."
        
    try:
        response = requests.post(
            pushover_url, 
            data={"token": pushover_token, "user": pushover_user, "message": text},
            timeout=10
        )
        response.raise_for_status()
        return "success"
    except requests.RequestException as e:
        return f"Push notification failed: {e}"



import pandas as pd
from session_manager import get_session_dir

ALLOWED_MODULES = {
    "math": math,
    "statistics": statistics,
    "random": random,
    "datetime": dt_module,
    "json": json,
    "re": re,
    "collections": collections,
    "itertools": itertools,
    "pandas": pd,
}

def safe_import(name, globals=None, locals=None, fromlist=(), level=0):
    """Runtime import filter allowing only pre-approved safe modules."""
    root_mod = name.split(".")[0]
    if root_mod in ALLOWED_MODULES:
        return __import__(name, globals, locals, fromlist, level)
    raise ImportError(f"Importing '{name}' is forbidden for security. Allowed: {list(ALLOWED_MODULES.keys())}")

# Safe sandbox definitions for Python REPL
SAFE_BUILTINS = {
    "__import__": safe_import,
    "abs": abs,
    "all": all,
    "any": any,
    "bin": bin,
    "bool": bool,
    "chr": chr,
    "dict": dict,
    "divmod": divmod,
    "enumerate": enumerate,
    "filter": filter,
    "float": float,
    "format": format,
    "frozenset": frozenset,
    "hex": hex,
    "int": int,
    "isinstance": isinstance,
    "issubclass": issubclass,
    "iter": iter,
    "len": len,
    "list": list,
    "map": map,
    "max": max,
    "min": min,
    "next": next,
    "oct": oct,
    "ord": ord,
    "pow": pow,
    "print": print,
    "range": range,
    "repr": repr,
    "reversed": reversed,
    "round": round,
    "set": set,
    "slice": slice,
    "sorted": sorted,
    "str": str,
    "sum": sum,
    "tuple": tuple,
    "zip": zip,
}

BLOCKED_NAMES = {
    "open", "eval", "exec", "compile", "getattr", "setattr", "delattr", 
    "globals", "locals", "vars", "input", "breakpoint", "exit", "quit",
    "help", "memoryview", "super"
}


def validate_python_code_ast(code_str: str) -> tuple[bool, str]:
    """Perform AST validation to ensure the code does not use dangerous constructs or imports."""
    try:
        tree = ast.parse(code_str)
    except SyntaxError as e:
        return False, f"SyntaxError: {e}"

    for node in ast.walk(tree):
        # Block private/dunder attribute access (__class__, __subclasses__, __globals__, etc.)
        if isinstance(node, ast.Attribute):
            if node.attr.startswith("_"):
                return False, f"Security Error: Access to private/dunder attribute '{node.attr}' is forbidden."
        
        # Block dunder names and dangerous built-in identifiers
        if isinstance(node, ast.Name):
            if node.id.startswith("__") and node.id.endswith("__"):
                return False, f"Security Error: Access to special identifier '{node.id}' is forbidden."
            if node.id in BLOCKED_NAMES:
                return False, f"Security Error: Call or access to '{node.id}' is forbidden."
        
        # Block imports outside of the strict whitelist
        if isinstance(node, ast.Import):
            for alias in node.names:
                root_mod = alias.name.split(".")[0]
                if root_mod not in ALLOWED_MODULES:
                    return False, f"Security Error: Importing '{alias.name}' is forbidden. Allowed modules: {list(ALLOWED_MODULES.keys())}."
        
        if isinstance(node, ast.ImportFrom):
            if node.module:
                root_mod = node.module.split(".")[0]
                if root_mod not in ALLOWED_MODULES:
                    return False, f"Security Error: Importing from '{node.module}' is forbidden. Allowed modules: {list(ALLOWED_MODULES.keys())}."

    return True, ""


def run_safe_python_code(query: str) -> str:
    """Execute Python code in a secure restricted namespace with AST inspection."""
    code = query.strip()
    if code.startswith("```python"):
        code = code[len("```python"):].strip()
    elif code.startswith("```"):
        code = code[len("```"):].strip()
    if code.endswith("```"):
        code = code[:-3].strip()

    is_safe, error_msg = validate_python_code_ast(code)
    if not is_safe:
        return error_msg

    namespace = {
        "__builtins__": SAFE_BUILTINS,
        **ALLOWED_MODULES
    }

    buffer = io.StringIO()
    try:
        with contextlib.redirect_stdout(buffer):
            exec(code, namespace, namespace)
        output = buffer.getvalue()
        return output if output else "Code executed successfully (no output produced. Remember to use print() if you need to see a result)."
    except Exception as e:
        return f"Execution Error: {type(e).__name__}: {e}"


def create_excel_report(data_input: str, filename: str = "enterprise_report.xlsx", session_id: str = "default") -> str:
    """Generate a styled, professional Excel report from structured JSON data."""
    try:
        if not filename.endswith(".xlsx"):
            filename = f"{filename}.xlsx"
        # Sanitize filename
        safe_filename = "".join(c for c in filename if c.isalnum() or c in "._-")
        if not safe_filename:
            safe_filename = "report.xlsx"

        session_dir = get_session_dir(session_id)
        filepath = os.path.join(session_dir, safe_filename)

        # Parse data input (accepts JSON string, list of dicts, or dict of tables)
        parsed_data = None
        try:
            parsed_data = json.loads(data_input)
        except Exception:
            # Try to handle lines of tabular text
            lines = [l.strip() for l in data_input.strip().split("\n") if l.strip()]
            if lines and ("|" in lines[0] or "," in lines[0]):
                sep = "|" if "|" in lines[0] else ","
                headers = [h.strip() for h in lines[0].split(sep) if h.strip()]
                rows = []
                for line in lines[1:]:
                    if "---" in line:
                        continue
                    row_vals = [v.strip() for v in line.split(sep) if v.strip()]
                    if row_vals:
                        rows.append(dict(zip(headers, row_vals)))
                parsed_data = rows

        if not parsed_data:
            return "Failed to generate Excel: Data input must be valid JSON array of objects or table structure."

        # Create DataFrame
        if isinstance(parsed_data, list):
            df = pd.DataFrame(parsed_data)
        elif isinstance(parsed_data, dict):
            df = pd.DataFrame(parsed_data)
        else:
            return "Failed to generate Excel: Unsupported data structure."

        # Write to Excel with openpyxl styling
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Executive Summary"
        ws.views.sheetView[0].showGridLines = True

        # Header styling
        header_fill = PatternFill(start_color="111827", end_color="111827", fill_type="solid")
        header_font = Font(name="Arial", size=11, bold=True, color="FFFFFF")
        header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        
        # Border styling
        thin_border = Border(
            left=Side(style='thin', color='E5E7EB'),
            right=Side(style='thin', color='E5E7EB'),
            top=Side(style='thin', color='E5E7EB'),
            bottom=Side(style='thin', color='E5E7EB')
        )
        alt_fill = PatternFill(start_color="F9FAFB", end_color="F9FAFB", fill_type="solid")

        # Write Headers
        headers = list(df.columns)
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=str(header))
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_align
            cell.border = thin_border
        ws.row_dimensions[1].height = 28

        # Write Data
        for row_idx, row_data in enumerate(df.values, 2):
            is_alt = (row_idx % 2 == 0)
            for col_idx, val in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=str(val) if val is not None else "")
                cell.font = Font(name="Arial", size=10, color="1F2937")
                cell.alignment = Alignment(vertical="center")
                cell.border = thin_border
                if is_alt:
                    cell.fill = alt_fill
            ws.row_dimensions[row_idx].height = 20

        # Auto-fit column widths
        for col in ws.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

        wb.save(filepath)
        return f"Successfully generated Excel workbook: '{safe_filename}' in session downloads. (Rows: {len(df)}, Columns: {len(headers)})"
    except Exception as e:
        return f"Error generating Excel report: {e}"


def create_executive_pdf(title: str, summary: str, sections_json: str, filename: str = "executive_brief.pdf", session_id: str = "default") -> str:
    """Generate a clean, high-impact Executive PDF Brief."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
        from reportlab.lib.units import inch

        if not filename.endswith(".pdf"):
            filename = f"{filename}.pdf"
        safe_filename = "".join(c for c in filename if c.isalnum() or c in "._-")
        if not safe_filename:
            safe_filename = "executive_brief.pdf"

        session_dir = get_session_dir(session_id)
        filepath = os.path.join(session_dir, safe_filename)

        doc = SimpleDocTemplate(
            filepath,
            pagesize=letter,
            rightMargin=40,
            leftMargin=40,
            topMargin=40,
            bottomMargin=40
        )
        styles = getSampleStyleSheet()

        # Custom Brand Styles
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=22,
            leading=26,
            textColor=colors.HexColor('#111827'),
            spaceAfter=6
        )
        
        meta_style = ParagraphStyle(
            'DocMeta',
            fontName='Helvetica',
            fontSize=9,
            leading=12,
            textColor=colors.HexColor('#6B7280'),
            spaceAfter=15
        )

        h2_style = ParagraphStyle(
            'SectionH2',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=14,
            leading=18,
            textColor=colors.HexColor('#111827'),
            spaceBefore=14,
            spaceAfter=8
        )

        body_style = ParagraphStyle(
            'BodyDark',
            parent=styles['BodyText'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#374151'),
            spaceAfter=6
        )

        summary_box_style = ParagraphStyle(
            'SummaryText',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10.5,
            leading=15,
            textColor=colors.HexColor('#065F46')
        )

        story = []
        
        # Title & Metadata
        story.append(Paragraph(title, title_style))
        current_date = datetime.now().strftime("%B %d, %Y - %H:%M")
        story.append(Paragraph(f"Sidekick Executive Intelligence Brief &bull; Generated: {current_date}", meta_style))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#E5E7EB'), spaceAfter=14))

        # Executive Summary Callout Box
        if summary:
            story.append(Paragraph("Executive Summary", h2_style))
            summary_p = Paragraph(f"<b>Key Takeaway:</b> {summary}", summary_box_style)
            summary_table = Table([[summary_p]], colWidths=[530])
            summary_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#F0FDF4')),
                ('BOX', (0, 0), (-1, -1), 1.5, colors.HexColor('#10B981')),
                ('PADDING', (0, 0), (-1, -1), 12),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ]))
            story.append(summary_table)
            story.append(Spacer(1, 14))

        # Parse Sections
        sections = []
        try:
            sections = json.loads(sections_json)
        except Exception:
            # Treat as plain text
            sections = [{"heading": "Detailed Findings", "text": str(sections_json)}]

        if isinstance(sections, list):
            for sec in sections:
                if isinstance(sec, dict):
                    h = sec.get("heading", "Analysis")
                    t = sec.get("text", "")
                    table_data = sec.get("table", None)

                    story.append(Paragraph(h, h2_style))
                    if t:
                        story.append(Paragraph(t.replace("\n", "<br/>"), body_style))
                    
                    if table_data and isinstance(table_data, list) and len(table_data) > 0:
                        # Build formatted table
                        t_rows = []
                        # If list of dicts:
                        if isinstance(table_data[0], dict):
                            cols = list(table_data[0].keys())
                            t_rows.append([Paragraph(f"<b>{c}</b>", body_style) for c in cols])
                            for r in table_data:
                                t_rows.append([Paragraph(str(r.get(c, '')), body_style) for c in cols])
                        # If list of lists:
                        elif isinstance(table_data[0], list):
                            for idx, row in enumerate(table_data):
                                is_head = (idx == 0)
                                t_rows.append([Paragraph(f"<b>{cell}</b>" if is_head else str(cell), body_style) for cell in row])

                        if t_rows:
                            col_w = 530 / max(len(t_rows[0]), 1)
                            pdf_table = Table(t_rows, colWidths=[col_w] * len(t_rows[0]))
                            pdf_table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F3F4F6')),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E5E7EB')),
                                ('PADDING', (0, 0), (-1, -1), 6),
                                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                            ]))
                            story.append(Spacer(1, 6))
                            story.append(pdf_table)
                            story.append(Spacer(1, 8))

        doc.build(story)
        return f"Successfully generated Executive PDF Brief: '{safe_filename}' in session downloads."
    except Exception as e:
        return f"Error generating PDF report: {e}"


def create_powerpoint_deck(title: str, subtitle: str, slides_json: str, filename: str = "executive_presentation.pptx", session_id: str = "default") -> str:
    """Generate a downloadable, high-impact 16:9 Executive PowerPoint Presentation (.pptx)."""
    try:
        session_dir = get_session_dir(session_id)
        safe_filename = os.path.basename(filename)
        if not safe_filename.endswith(".pptx"):
            safe_filename += ".pptx"
        output_path = os.path.join(session_dir, safe_filename)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # Parse slides JSON
        if isinstance(slides_json, str):
            try:
                slides_data = json.loads(slides_json)
            except Exception:
                slides_data = []
        else:
            slides_data = slides_json or []

        # 1. Title Slide (Executive Dark Slate Theme)
        title_slide = prs.slides.add_slide(blank_layout)
        bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
        bg.line.fill.background()

        # Accent vertical bar
        accent = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.0), Inches(0.12), Inches(3.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(16, 185, 129) # Emerald 500
        accent.line.fill.background()

        # Title Box
        title_box = title_slide.shapes.add_textbox(Inches(1.6), Inches(2.0), Inches(10.5), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle Box
        sub_box = title_slide.shapes.add_textbox(Inches(1.6), Inches(4.0), Inches(10.5), Inches(1.2))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        p_sub = sub_tf.paragraphs[0]
        p_sub.text = subtitle or "Strategic Executive Presentation"
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(148, 163, 184) # Slate 400

        # Footer badge
        foot_box = title_slide.shapes.add_textbox(Inches(1.6), Inches(5.6), Inches(10.5), Inches(0.6))
        foot_tf = foot_box.text_frame
        p_foot = foot_tf.paragraphs[0]
        cur_date = datetime.now().strftime("%B %d, %Y")
        p_foot.text = f"⚡ Sidekick Executive Intelligence • Generated {cur_date}"
        p_foot.font.size = Pt(13)
        p_foot.font.bold = True
        p_foot.font.color.rgb = RGBColor(16, 185, 129)

        # 2. Content Slides (Corporate Clean Light Theme)
        for s_idx, slide_info in enumerate(slides_data):
            c_slide = prs.slides.add_slide(blank_layout)
            slide_title = slide_info.get("title", f"Section {s_idx + 1}")
            category = slide_info.get("category", "EXECUTIVE BRIEFING")

            # Header Banner
            head_box = c_slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.95))
            h_tf = head_box.text_frame
            h_tf.word_wrap = True
            
            p_cat = h_tf.paragraphs[0]
            p_cat.text = category.upper()
            p_cat.font.size = Pt(10.5)
            p_cat.font.bold = True
            p_cat.font.color.rgb = RGBColor(5, 150, 105) # Emerald 600

            p_h = h_tf.add_paragraph()
            p_h.text = slide_title
            p_h.font.size = Pt(24)
            p_h.font.bold = True
            p_h.font.color.rgb = RGBColor(15, 23, 42)

            # Divider line
            div = c_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.42), Inches(11.733), Inches(0.015))
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(226, 232, 240)
            div.line.fill.background()

            # Render Corporate Cards (e.g. 2, 3, or 4 columns)
            cards = slide_info.get("cards", [])
            if cards:
                num_cards = min(len(cards), 4)
                total_w = 11.733
                spacing = 0.3
                card_w = (total_w - (num_cards - 1) * spacing) / num_cards
                card_h = Inches(4.9)
                top_y = Inches(1.55)

                for c_i, card in enumerate(cards[:num_cards]):
                    left_x = Inches(0.8 + c_i * (card_w + spacing))
                    
                    # Card Outer Container
                    c_box = c_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, Inches(card_w), card_h)
                    c_box.fill.solid()
                    c_box.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
                    c_box.line.color.rgb = RGBColor(226, 232, 240) # Slate 200
                    c_box.line.width = Pt(1.0)

                    # Card Header Bar (Crisp corporate header strip)
                    c_head = c_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, Inches(card_w), Inches(0.85))
                    c_head.fill.solid()
                    c_head.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                    c_head.line.color.rgb = RGBColor(226, 232, 240)
                    c_head.line.width = Pt(1.0)

                    c_tf = c_head.text_frame
                    c_tf.word_wrap = True
                    c_tf.margin_left = Inches(0.18)
                    c_tf.margin_top = Inches(0.12)

                    p_ctitle = c_tf.paragraphs[0]
                    p_ctitle.text = card.get("title", f"Item {c_i + 1}")
                    p_ctitle.font.size = Pt(16)
                    p_ctitle.font.bold = True
                    p_ctitle.font.color.rgb = RGBColor(15, 23, 42)

                    if card.get("subtitle"):
                        p_csub = c_tf.add_paragraph()
                        p_csub.text = card.get("subtitle")
                        p_csub.font.size = Pt(11)
                        p_csub.font.color.rgb = RGBColor(100, 116, 139)

                    # Card Body Text / Bullets
                    b_box = c_slide.shapes.add_textbox(left_x, top_y + Inches(0.95), Inches(card_w), Inches(3.1))
                    b_tf = b_box.text_frame
                    b_tf.word_wrap = True
                    b_tf.margin_left = Inches(0.18)
                    b_tf.margin_right = Inches(0.18)
                    b_tf.margin_top = Inches(0.1)

                    points = card.get("points", [])
                    for p_idx, point in enumerate(points):
                        p_pt = b_tf.paragraphs[0] if p_idx == 0 else b_tf.add_paragraph()
                        p_pt.text = f"•  {point}"
                        p_pt.font.size = Pt(12)
                        p_pt.font.color.rgb = RGBColor(30, 41, 59)
                        p_pt.space_after = Pt(6)

                    # Card Bottom Highlight Badge
                    if card.get("highlight"):
                        hl_box = c_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x + Inches(0.15), top_y + Inches(4.2), Inches(card_w - 0.3), Inches(0.55))
                        hl_box.fill.solid()
                        hl_box.fill.fore_color.rgb = RGBColor(236, 253, 245) # Emerald 50
                        hl_box.line.color.rgb = RGBColor(110, 231, 183) # Emerald 300
                        hl_box.line.width = Pt(0.75)

                        hl_tf = hl_box.text_frame
                        hl_tf.word_wrap = True
                        hl_tf.margin_left = Inches(0.1)
                        hl_tf.margin_top = Inches(0.08)
                        p_hl = hl_tf.paragraphs[0]
                        p_hl.text = f"★ {card.get('highlight')}"
                        p_hl.font.size = Pt(11)
                        p_hl.font.bold = True
                        p_hl.font.color.rgb = RGBColor(4, 120, 87)

            # Or render Bullets & Text
            elif slide_info.get("bullets"):
                b_box = c_slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(4.5))
                b_tf = b_box.text_frame
                b_tf.word_wrap = True
                
                bullets = slide_info.get("bullets", [])
                for b_i, bullet in enumerate(bullets):
                    p_b = b_tf.paragraphs[0] if b_i == 0 else b_tf.add_paragraph()
                    p_b.text = f"•  {bullet}"
                    p_b.font.size = Pt(15)
                    p_b.font.color.rgb = RGBColor(30, 41, 59)
                    p_b.space_after = Pt(10)

            # Bottom Strategic Takeaway Bar if present
            if slide_info.get("takeaway"):
                t_box = c_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.55), Inches(11.733), Inches(0.58))
                t_box.fill.solid()
                t_box.fill.fore_color.rgb = RGBColor(240, 253, 244) # Emerald 50
                t_box.line.color.rgb = RGBColor(16, 185, 129) # Emerald 500
                t_box.line.width = Pt(1.0)

                t_tf = t_box.text_frame
                t_tf.word_wrap = True
                t_tf.margin_left = Inches(0.15)
                t_tf.margin_top = Inches(0.08)
                p_t = t_tf.paragraphs[0]
                p_t.text = f"💡 STRATEGIC TAKEAWAY: {slide_info.get('takeaway')}"
                p_t.font.size = Pt(11.5)
                p_t.font.bold = True
                p_t.font.color.rgb = RGBColor(6, 95, 70)

        prs.save(output_path)
        return f"Successfully generated Executive PowerPoint Presentation: '{safe_filename}' in session downloads. (Slides: {len(slides_data) + 1})"
    except Exception as e:
        return f"Error generating PowerPoint presentation: {e}"


class ExcelReportInput(BaseModel):
    data_json: str = Field(description="JSON array of objects representing rows (e.g. '[{\"Company\": \"A\", \"Pricing\": \"$50\", \"Score\": 95}]') or a pipe-separated table string.")
    filename: str = Field(default="market_analysis.xlsx", description="Filename for the Excel file, ending with .xlsx.")


class ExecutivePdfInput(BaseModel):
    title: str = Field(description="Title of the Executive Brief.")
    summary: str = Field(description="Executive Summary takeaway paragraph.")
    sections_json: str = Field(default="[]", description="JSON array of objects with keys 'heading', 'text', and optional 'table'.")
    filename: str = Field(default="executive_brief.pdf", description="Filename for the PDF file, ending with .pdf.")


class PowerPointDeckInput(BaseModel):
    title: str = Field(description="Title of the Executive PowerPoint Presentation.")
    subtitle: str = Field(default="Strategic Executive Briefing", description="Subtitle of the presentation.")
    slides_json: str = Field(description="JSON array of slide objects. Each object can have 'title', 'category', 'cards' (list of card objects with 'title', 'subtitle', 'points', 'highlight'), 'bullets' (list of bullet strings), and optional 'takeaway' (string).")
    filename: str = Field(default="executive_presentation.pptx", description="Filename for the PowerPoint presentation, ending with .pptx.")


async def other_tools(session_id: str = "default"):
    push_tool = Tool(
        name="send_push_notification", 
        func=push, 
        description="Use this tool when you want to send a push notification alert when a task is completed."
    )
    
    tool_search = Tool(
        name="search",
        func=serp.run,
        description="Use this tool when you want to get current search results from Google."
    )

    wikipedia = WikipediaAPIWrapper()
    wiki_tool = WikipediaQueryRun(api_wrapper=wikipedia)

    python_repl = Tool(
        name="python_repl",
        func=run_safe_python_code,
        description="A sandboxed Python REPL with pandas, math, statistics, and regex. Use this to execute data processing, statistical calculations, or algorithms. Use print(...) to output results."
    )

    def excel_tool_fn(data_json: str, filename: str = "market_analysis.xlsx") -> str:
        return create_excel_report(data_json, filename, session_id=session_id)

    excel_tool = StructuredTool.from_function(
        func=excel_tool_fn,
        name="generate_excel_report",
        description="Generate a downloadable, professional Microsoft Excel (.xlsx) spreadsheet with auto-styled headers and alternating rows.",
        args_schema=ExcelReportInput
    )

    def pdf_tool_fn(title: str, summary: str, sections_json: str = "[]", filename: str = "executive_brief.pdf") -> str:
        return create_executive_pdf(title, summary, sections_json, filename, session_id=session_id)

    pdf_tool = StructuredTool.from_function(
        func=pdf_tool_fn,
        name="generate_executive_pdf",
        description="Generate a downloadable, high-impact Executive PDF Brief complete with branded header, executive summary callout box, structured sub-sections, and tables.",
        args_schema=ExecutivePdfInput
    )

    def pptx_tool_fn(title: str, subtitle: str = "Strategic Executive Briefing", slides_json: str = "[]", filename: str = "executive_presentation.pptx") -> str:
        return create_powerpoint_deck(title, subtitle, slides_json, filename, session_id=session_id)

    pptx_tool = StructuredTool.from_function(
        func=pptx_tool_fn,
        name="generate_powerpoint_presentation",
        description="Generate a downloadable, high-impact 16:9 Executive PowerPoint Presentation (.pptx) with modern title slide, multi-column card comparisons, bullet slides, and strategic takeaways.",
        args_schema=PowerPointDeckInput
    )
    
    return [push_tool, tool_search, python_repl, wiki_tool, excel_tool, pdf_tool, pptx_tool]



